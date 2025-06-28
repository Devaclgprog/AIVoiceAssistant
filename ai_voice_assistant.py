#final one 
import streamlit as st
import tempfile
from faster_whisper import WhisperModel
from fpdf import FPDF
import google.generativeai as genai
from datetime import datetime
import pandas as pd
import PyPDF2
import os
from io import StringIO
import random
from pptx import Presentation
from pptx.util import Pt
import base64
import sys
from dotenv import load_dotenv


load_dotenv()
st.set_page_config(
    page_title="Voice AI Assistant",
    page_icon="🎤",
    layout="centered",
    initial_sidebar_state="collapsed"
)


genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
model = genai.GenerativeModel('gemini-2.0-flash')


def text_to_pdf(text, filename):
    """Convert text content to a PDF file"""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, txt=text)
    pdf.output(filename)
    return filename

def get_binary_file_downloader_html(bin_file, file_label='File'):
    """Generate a download link for binary files"""
    with open(bin_file, 'rb') as f:
        data = f.read()
    bin_str = base64.b64encode(data).decode()
    href = f'<a href="data:application/octet-stream;base64,{bin_str}" download="{os.path.basename(bin_file)}">{file_label}</a>'
    return href


if 'transcription' not in st.session_state:
    st.session_state.update({
        "page": "main",
        "transcription": "",
        "original_transcription": "",
        "summary": "",
        "audio_data": None,
        "audio_path": "",
        "messages": [],
        "ppt_path": None,
        "summary_path": None,
        "ppt_title": "",
        "ppt_headings": "",
        "transcript_edited": False,
        "ppt_edited": False,
        "uploaded_files": [],
        "file_contents": {},
        "previous_questions": {}
    })


def extract_text_from_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    return "\n".join(page.extract_text() for page in pdf_reader.pages)

def extract_text_from_csv(file):
    return pd.read_csv(file).to_string()

def process_uploaded_files():
    for file in st.session_state.uploaded_files:
        if file.name not in st.session_state.file_contents:
            if file.name.endswith('.pdf'):
                st.session_state.file_contents[file.name] = f"PDF Content:\n{extract_text_from_pdf(file)}"
            elif file.name.endswith('.csv'):
                st.session_state.file_contents[file.name] = f"CSV Content:\n{extract_text_from_csv(file)}"

def transcribe_audio():
    with st.spinner("Transcribing audio..."):
        model = WhisperModel("base", device="cpu", compute_type="int8")
        segments, _ = model.transcribe(st.session_state.audio_path)
        transcription = "\n".join(f"[{seg.start:.2f}s] {seg.text}" for seg in segments)
        st.session_state.original_transcription = transcription
        return transcription

def generate_summary():
    try:
        transcript = st.session_state.transcription if st.session_state.transcript_edited else st.session_state.original_transcription
        response = model.generate_content(f"""
            Create a professional summary from this transcript:
            {transcript}
            Include key points, action items, and recommendations.
            Use markdown formatting with headings and bullet points.
        """)
        return response.text
    except Exception as e:
        return f"❌ Error: {str(e)}"

def create_presentation(title, headings):
    try:
        prs = Presentation()
        
        
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
        
       
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = f"Generated {datetime.now().strftime('%d %b %Y %H:%M')}"
            for paragraph in subtitle.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
        
        for heading in [h.strip() for h in headings.split('\n') if h.strip()]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            
            title_shape = slide.shapes.title
            title_shape.text = heading
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
            
            
            if len(slide.placeholders) > 1:
                content = slide.placeholders[1]
                transcript = st.session_state.transcription if st.session_state.transcript_edited else st.session_state.original_transcription
                response = model.generate_content(f"Create 3-5 bullet points about '{heading}' using: {transcript}")
                
                text_frame = content.text_frame
                text_frame.word_wrap = True
                for point in [p.strip() for p in response.text.split('\n') if p.strip()]:
                    p = text_frame.add_paragraph()
                    p.text = point.replace("- ", "").strip()
                    p.level = 0
                    p.font.size = Pt(9)
                    p.space_after = Pt(9)
        
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_file:
            prs.save(tmp_file.name)
            return tmp_file.name
            
    except Exception as e:
        st.error(f"PPT Error: {str(e)}")
        return None

def chat_response(prompt):
    
    if len(st.session_state.messages) >= 2:
        last_user_msg = next((msg["content"] for msg in reversed(st.session_state.messages) if msg["role"] == "user"), "")
        if prompt.lower().strip() == last_user_msg.lower().strip():
            variations = [
                "I believe I already answered this, but to recap:",
                "As I mentioned earlier:",
                "Let me rephrase my previous response:",
                "To reiterate what I shared before:",
                "Just to summarize again:"
            ]
            recap_intro = random.choice(variations)
    
    transcript = st.session_state.transcription if st.session_state.transcript_edited else st.session_state.original_transcription
    
   
    file_context = ""
    if st.session_state.file_contents:
        file_context = "\n\nAdditional Reference Documents:\n" + "\n".join(
            f"=== {filename} ===\n{content}\n"
            for filename, content in st.session_state.file_contents.items()
        )
    
    
    system_prompt = f"""
    ROLE: Professional Human-like Assistant
    CONTEXT: {transcript}
    {file_context}
    CONVERSATION HISTORY: {st.session_state.messages[-6:]}
    
    GUIDELINES:
    1. Respond naturally like a human colleague would
    2. Maintain continuity - reference previous exchanges naturally
    3. For repeated questions, acknowledge and vary responses
    4. Keep tone professional but friendly
    5. Use natural language variations
    6. When appropriate, add brief conversational elements
    7. Format responses clearly with markdown when helpful
    1. Answer based on both the audio content and any uploaded documents
    2. If question is unrelated to available content, respond politely
    3. Maintain conversational flow
    4. You are a chatbot named "Chat with your voice and documents"
    5. Use markdown formatting when appropriate
    6. Keep responses concise but helpful
    7. Reference specific documents when relevant
    8. Never mention you're an AI or language model
    9. Maintain friendly, professional tone
    10. Acknowledge previous interactions when relevant
    CURRENT QUESTION: {prompt}
    """
    
    try:
        response = model.generate_content(system_prompt)
        text = response.text
        
        
        if 'recap_intro' in locals():
            text = f"{recap_intro}\n\n{text}"
        
        if random.random() < 0.2:
            conversational_addons = [
                "\n\nLet me know if you'd like me to elaborate on any part.",
                "\n\nDoes this help answer your question?",
                "\n\nI'm happy to discuss this further if needed.",
                "\n\nWhat else would you like to know about this?",
                "\n\nWould you like me to approach this from a different angle?"
            ]
            text += random.choice(conversational_addons)
            
        return text
        
    except Exception as e:
        return f"⚠️ Error: {str(e)}"


def main_page():
    st.title("🎤 Voice AI Assistant")
    
    
    col1, col2 = st.columns([3, 1])
    with col1:
        if st.session_state.audio_data is None:
            st.session_state.audio_data = st.audio_input("Record Audio", key="recorder")
        else:
            st.audio(st.session_state.audio_data, format="audio/wav")
    
    with col2:
        if st.button("🔄 New Session"):
            st.session_state.clear()
            st.rerun()
    
    
    if st.session_state.audio_data and not st.session_state.transcription:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as tmp_audio:
            tmp_audio.write(st.session_state.audio_data.read())
            st.session_state.audio_path = tmp_audio.name
        st.session_state.transcription = transcribe_audio()
    
    if st.session_state.transcription:
        with st.expander("📜 Transcript"):
            edited = st.text_area("Edit", st.session_state.transcription, height=300)
            if edited != st.session_state.transcription:
                st.session_state.transcription = edited
                st.session_state.transcript_edited = True
            
          
            col1, col2 = st.columns(2)
            with col1:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_pdf:
                    text_to_pdf(st.session_state.transcription, tmp_pdf.name)
                    with open(tmp_pdf.name,"rb") as f:
                        st.download_button(
                            "⬇️ Current Transcript (PDF)",
                            data =f,
                            file_name ="current_transcript.pdf",
                            mime = "application/pdf"
                        )
            with col2:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_pdf:
                    text_to_pdf(st.session_state.original_transcription, tmp_pdf.name)
                    with open(tmp_pdf.name,"rb") as f:
                        st.download_button(
                            "⬇️ Original Transcript (PDF)",
                            data =f,
                            file_name ="original_transcript.pdf",
                            mime = "application/pdf"
                        )
                    
     
        cols = st.columns(3)
        cols[0].button("📄 Summary", on_click=lambda: st.session_state.update({"page": "summary"}))
        cols[1].button("📊 PPT", on_click=lambda: st.session_state.update({"page": "ppt"}))
        cols[2].button("💬 Chat", on_click=lambda: st.session_state.update({"page": "chat"}))

def summary_page():
    st.title("📄 Summary")
    if st.button("← Back"): 
        st.session_state.page = "main"
    
    if not st.session_state.summary:
        st.session_state.summary = generate_summary()
    
    st.markdown(st.session_state.summary)
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_pdf:
        text_to_pdf(st.session_state.summary, tmp_pdf.name)
        with open(tmp_pdf.name,"rb") as f:
            st.download_button(
                            "⬇️ Download Summary(PDF)",
                            data =f,
                            file_name ="summary.pdf",
                            mime = "application/pdf"
                        )

def ppt_page():
    st.title("📊 PPT Generator")
    if st.button("← Back"): 
        st.session_state.page = "main"
    
    try:
        from pptx import Presentation
        pptx_available = True
    except ImportError:
        pptx_available = False
        st.error("PPT generation requires python-pptx package. Please install with: pip install python-pptx")
        return
    
    with st.form("ppt_form"):
        title = st.text_input("Title", "Voice Analysis Report")
        headings = st.text_area("Headings (one per line)", "Introduction\nFindings\nRecommendations")
        if st.form_submit_button("Generate"):
            st.session_state.ppt_path = create_presentation(title, headings)
    
    if st.session_state.ppt_path:
        st.success("✅ Done!")
        st.download_button("⬇️ Download", open(st.session_state.ppt_path, "rb"), "presentation.pptx")

def chat_page():
    st.title("💬 Chat")
    if st.button("← Back"): 
        st.session_state.page = "main"
    
 
    with st.expander("📁 Upload Files"):
        uploaded = st.file_uploader("PDF/CSV", ["pdf", "csv"], True)
        if uploaded and st.button("Process"):
            st.session_state.uploaded_files = uploaded
            process_uploaded_files()
    
   
    for msg in st.session_state.messages:
        st.chat_message(msg["role"]).write(msg["content"])
    
    
    if prompt := st.chat_input("Ask about the recording"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.spinner("Thinking..."):
            st.session_state.messages.append({"role": "assistant", "content": chat_response(prompt)})
        st.rerun()


if st.session_state.page == "main":
    main_page()
elif st.session_state.page == "summary":
    summary_page()
elif st.session_state.page == "ppt":
    ppt_page()
elif st.session_state.page == "chat":
    chat_page()
else:
    st.session_state.page = "main"
    st.rerun()
